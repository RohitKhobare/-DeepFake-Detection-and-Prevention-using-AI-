"""
Comprehensive File Content Analyzer
Analyzes complete file content, metadata, and determines AI vs human creation
"""

import os
import json
import hashlib
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Any
import tempfile
import logging
from datetime import datetime
import mimetypes
try:
    import magic
    MAGIC_AVAILABLE = True
except ImportError:
    MAGIC_AVAILABLE = False
    print("⚠️ python-magic not available, using fallback file type detection")
import exifread
from PIL import Image, ExifTags
import fitz  # PyMuPDF for PDF analysis
import docx
import pptx
import openpyxl
from moviepy import VideoFileClip
import librosa
import soundfile as sf
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.naive_bayes import MultinomialNB
import pickle
import numpy as np

logger = logging.getLogger(__name__)


class ContentAnalyzer:
    """Analyzes file content comprehensively"""

    def __init__(self):
        self.supported_formats = {
            # Images
            'image': ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp', '.ico', '.svg', '.raw', '.heic', '.gif'],
            # Videos
            'video': ['.mp4', '.avi', '.mov', '.mkv', '.webm', '.flv', '.wmv', '.m4v', '.mpg', '.3gp', '.asf', '.rm'],
            # Audio
            'audio': ['.mp3', '.wav', '.flac', '.aac', '.m4a', '.wma', '.ogg', '.opus', '.aiff', '.au'],
            # Documents
            'document': ['.pdf', '.docx', '.pptx', '.xlsx', '.txt', '.rtf', '.odt', '.ods', '.odp'],
            # Archives
            'archive': ['.zip', '.rar', '.7z', '.tar', '.gz', '.bz2'],
            # Other
            'other': ['.json', '.xml', '.csv', '.html', '.css', '.js', '.py', '.cpp', '.java', '.cs']
        }

        # Load AI detection models (placeholder - would be trained models)
        self.ai_detection_models = self._load_ai_detection_models()

    def _load_ai_detection_models(self) -> Dict:
        """Load pre-trained AI detection models"""
        # Placeholder for actual model loading
        # In production, these would be trained models for different content types
        return {
            'text': None,  # Would be a language model for text analysis
            'image': None,  # Would be CNN-based model for image analysis
            'audio': None,  # Would be audio analysis model
            'video': None   # Would be video analysis model
        }

    def analyze_file_comprehensive(self, file_path: str) -> Dict[str, Any]:
        """Perform comprehensive file analysis"""
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                return {'error': 'File does not exist'}

            # Basic file information
            file_info = self._get_basic_file_info(file_path)

            # Content type detection
            content_type = self._detect_content_type(file_path)

            # Content analysis based on type
            content_analysis = self._analyze_content_by_type(file_path, content_type)

            # Metadata extraction
            metadata = self._extract_metadata(file_path, content_type)

            # AI vs Human detection
            ai_human_analysis = self._analyze_ai_vs_human(file_path, content_type, content_analysis)

            # File summary
            summary = self._generate_file_summary(file_info, content_analysis, metadata, ai_human_analysis)

            return {
                'file_info': file_info,
                'content_type': content_type,
                'content_analysis': content_analysis,
                'metadata': metadata,
                'ai_human_analysis': ai_human_analysis,
                'summary': summary,
                'analysis_timestamp': datetime.now().isoformat(),
                'analyzer_version': '2.0.0'
            }

        except Exception as e:
            logger.error(f"Comprehensive analysis error for {file_path}: {e}")
            return {
                'error': str(e),
                'file_path': str(file_path),
                'analysis_timestamp': datetime.now().isoformat()
            }

    def _get_basic_file_info(self, file_path: Path) -> Dict:
        """Get basic file information"""
        try:
            stat = file_path.stat()
            return {
                'file_path': str(file_path),
                'file_name': file_path.name,
                'file_size': stat.st_size,
                'file_size_human': self._format_file_size(stat.st_size),
                'created_time': datetime.fromtimestamp(stat.st_ctime).isoformat(),
                'modified_time': datetime.fromtimestamp(stat.st_mtime).isoformat(),
                'accessed_time': datetime.fromtimestamp(stat.st_atime).isoformat(),
                'file_extension': file_path.suffix.lower(),
                'mime_type': mimetypes.guess_type(str(file_path))[0],
                'magic_type': magic.from_file(str(file_path), mime=True) if MAGIC_AVAILABLE and hasattr(magic, 'from_file') else None
            }
        except Exception as e:
            return {'error': f'Cannot get file info: {e}'}

    def _detect_content_type(self, file_path: Path) -> str:
        """Detect the content type of the file"""
        ext = file_path.suffix.lower()

        for category, extensions in self.supported_formats.items():
            if ext in extensions:
                return category

        # Fallback to mime type detection
        mime_type = mimetypes.guess_type(str(file_path))[0]
        if mime_type:
            if mime_type.startswith('image/'):
                return 'image'
            elif mime_type.startswith('video/'):
                return 'video'
            elif mime_type.startswith('audio/'):
                return 'audio'
            elif mime_type in ['application/pdf', 'application/msword', 'application/vnd.openxmlformats-officedocument']:
                return 'document'

        return 'other'

    def _analyze_content_by_type(self, file_path: Path, content_type: str) -> Dict:
        """Analyze content based on file type"""
        try:
            if content_type == 'image':
                return self._analyze_image_content(file_path)
            elif content_type == 'video':
                return self._analyze_video_content(file_path)
            elif content_type == 'audio':
                return self._analyze_audio_content(file_path)
            elif content_type == 'document':
                return self._analyze_document_content(file_path)
            elif content_type == 'archive':
                return self._analyze_archive_content(file_path)
            else:
                return self._analyze_generic_content(file_path)
        except Exception as e:
            return {'error': f'Content analysis failed: {e}'}

    def _analyze_image_content(self, file_path: Path) -> Dict:
        """Comprehensive image content analysis"""
        try:
            with Image.open(file_path) as img:
                # Basic properties
                analysis = {
                    'dimensions': img.size,
                    'mode': img.mode,
                    'format': img.format,
                    'has_transparency': img.mode in ('RGBA', 'LA', 'P'),
                    'color_depth': self._get_color_depth(img),
                    'aspect_ratio': img.size[0] / img.size[1] if img.size[1] > 0 else 0
                }

                # Color analysis
                analysis['color_analysis'] = self._analyze_image_colors(img)

                # Compression analysis
                analysis['compression_analysis'] = self._analyze_image_compression(img, file_path)

                # Quality analysis
                analysis['quality_analysis'] = self._analyze_image_quality(img)

                # Pattern analysis
                analysis['pattern_analysis'] = self._analyze_image_patterns(img)

                # Metadata analysis
                analysis['embedded_metadata'] = self._extract_image_metadata(img)

                return analysis

        except Exception as e:
            return {'error': f'Image analysis failed: {e}'}

    def _analyze_video_content(self, file_path: Path) -> Dict:
        """Comprehensive video content analysis"""
        try:
            clip = VideoFileClip(str(file_path))

            analysis = {
                'duration': clip.duration,
                'fps': clip.fps,
                'size': clip.size,
                'aspect_ratio': clip.size[0] / clip.size[1] if clip.size[1] > 0 else 0,
                'audio_present': clip.audio is not None,
                'frame_count': int(clip.duration * clip.fps) if clip.duration and clip.fps else 0
            }

            # Audio analysis if present
            if clip.audio:
                analysis['audio_analysis'] = self._analyze_audio_from_video(clip.audio)

            # Frame analysis (sample frames)
            analysis['frame_analysis'] = self._analyze_video_frames(clip)

            # Quality analysis
            analysis['quality_analysis'] = self._analyze_video_quality(clip)

            # Content analysis
            analysis['content_analysis'] = self._analyze_video_content_features(clip)

            clip.close()
            return analysis

        except Exception as e:
            return {'error': f'Video analysis failed: {e}'}

    def _analyze_audio_content(self, file_path: Path) -> Dict:
        """Comprehensive audio content analysis"""
        try:
            # Load audio
            y, sr = librosa.load(str(file_path), sr=None)

            analysis = {
                'duration': len(y) / sr,
                'sample_rate': sr,
                'channels': 1 if y.ndim == 1 else y.shape[0],
                'total_samples': len(y)
            }

            # Quality analysis
            analysis['quality_analysis'] = self._analyze_audio_quality(y, sr)

            # Content analysis
            analysis['content_analysis'] = self._analyze_audio_content(y, sr)

            # Technical analysis
            analysis['technical_analysis'] = self._analyze_audio_technical(y, sr)

            return analysis

        except Exception as e:
            return {'error': f'Audio analysis failed: {e}'}

    def _analyze_document_content(self, file_path: Path) -> Dict:
        """Comprehensive document content analysis"""
        try:
            ext = file_path.suffix.lower()

            if ext == '.pdf':
                return self._analyze_pdf_content(file_path)
            elif ext == '.docx':
                return self._analyze_docx_content(file_path)
            elif ext == '.pptx':
                return self._analyze_pptx_content(file_path)
            elif ext == '.xlsx':
                return self._analyze_xlsx_content(file_path)
            elif ext == '.txt':
                return self._analyze_text_content(file_path)
            else:
                return {'error': 'Unsupported document format'}

        except Exception as e:
            return {'error': f'Document analysis failed: {e}'}

    def _analyze_pdf_content(self, file_path: Path) -> Dict:
        """Analyze PDF content"""
        try:
            doc = fitz.open(str(file_path))

            analysis = {
                'page_count': len(doc),
                'metadata': dict(doc.metadata),
                'pages': []
            }

            # Analyze each page
            for i, page in enumerate(doc):
                page_analysis = {
                    'page_number': i + 1,
                    'text_length': len(page.get_text()),
                    'image_count': len(page.get_images(full=True)),
                    'link_count': len(page.get_links()),
                    'annotation_count': len(page.annotations())
                }
                analysis['pages'].append(page_analysis)

            # Overall statistics
            analysis['total_text_length'] = sum(p['text_length'] for p in analysis['pages'])
            analysis['total_images'] = sum(p['image_count'] for p in analysis['pages'])
            analysis['total_links'] = sum(p['link_count'] for p in analysis['pages'])

            doc.close()
            return analysis

        except Exception as e:
            return {'error': f'PDF analysis failed: {e}'}

    def _analyze_docx_content(self, file_path: Path) -> Dict:
        """Analyze DOCX content"""
        try:
            doc = docx.Document(str(file_path))

            analysis = {
                'paragraph_count': len(doc.paragraphs),
                'table_count': len(doc.tables),
                'section_count': len(doc.sections),
                'total_text_length': sum(len(p.text) for p in doc.paragraphs),
                'properties': {
                    'author': doc.core_properties.author,
                    'created': doc.core_properties.created.isoformat() if doc.core_properties.created else None,
                    'modified': doc.core_properties.modified.isoformat() if doc.core_properties.modified else None,
                    'title': doc.core_properties.title
                }
            }

            return analysis

        except Exception as e:
            return {'error': f'DOCX analysis failed: {e}'}

    def _analyze_text_content(self, file_path: Path) -> Dict:
        """Analyze plain text content"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()

            analysis = {
                'character_count': len(content),
                'line_count': len(content.split('\n')),
                'word_count': len(content.split()),
                'paragraph_count': len([p for p in content.split('\n\n') if p.strip()]),
                'language_detection': self._detect_text_language(content),
                'readability_score': self._calculate_readability(content),
                'sentiment_analysis': self._analyze_text_sentiment(content)
            }

            return analysis

        except Exception as e:
            return {'error': f'Text analysis failed: {e}'}

    def _extract_metadata(self, file_path: Path, content_type: str) -> Dict:
        """Extract comprehensive metadata"""
        metadata = {}

        try:
            if content_type == 'image':
                metadata.update(self._extract_image_metadata_file(file_path))
            elif content_type == 'video':
                metadata.update(self._extract_video_metadata(file_path))
            elif content_type == 'audio':
                metadata.update(self._extract_audio_metadata(file_path))
            elif content_type == 'document':
                metadata.update(self._extract_document_metadata(file_path))

            # File system metadata
            metadata['file_system'] = self._get_file_system_metadata(file_path)

            # Hash for integrity checking
            metadata['hashes'] = self._calculate_file_hashes(file_path)

        except Exception as e:
            metadata['error'] = f'Metadata extraction failed: {e}'

        return metadata

    def _analyze_ai_vs_human(self, file_path: Path, content_type: str, content_analysis: Dict) -> Dict:
        """Analyze whether content was likely created by AI or human"""
        try:
            analysis = {
                'content_type': content_type,
                'ai_probability': 0.5,
                'human_probability': 0.5,
                'confidence': 'medium',
                'indicators': []
            }

            if content_type == 'image':
                analysis.update(self._analyze_image_ai_vs_human(content_analysis))
            elif content_type == 'video':
                analysis.update(self._analyze_video_ai_vs_human(content_analysis))
            elif content_type == 'audio':
                analysis.update(self._analyze_audio_ai_vs_human(content_analysis))
            elif content_type == 'document':
                analysis.update(self._analyze_document_ai_vs_human(content_analysis))

            # Overall assessment
            if analysis['ai_probability'] > 0.7:
                analysis['assessment'] = 'Likely AI-generated'
            elif analysis['human_probability'] > 0.7:
                analysis['assessment'] = 'Likely human-created'
            else:
                analysis['assessment'] = 'Uncertain origin'

            return analysis

        except Exception as e:
            return {'error': f'AI vs Human analysis failed: {e}'}

    def _generate_file_summary(self, file_info: Dict, content_analysis: Dict,
                             metadata: Dict, ai_human_analysis: Dict) -> Dict:
        """Generate comprehensive file summary"""
        try:
            summary = {
                'file_name': file_info.get('file_name', 'Unknown'),
                'file_type': file_info.get('file_extension', 'Unknown'),
                'file_size': file_info.get('file_size_human', 'Unknown'),
                'content_summary': self._summarize_content(content_analysis),
                'key_findings': self._extract_key_findings(content_analysis, metadata, ai_human_analysis),
                'risk_assessment': self._assess_risks(ai_human_analysis),
                'recommendations': self._generate_recommendations(ai_human_analysis)
            }

            return summary

        except Exception as e:
            return {'error': f'Summary generation failed: {e}'}

    # Helper methods for detailed analysis
    def _analyze_image_colors(self, img: Image.Image) -> Dict:
        """Analyze image color properties"""
        try:
            img_array = np.array(img)

            if img.mode == 'RGB':
                r, g, b = img_array[:, :, 0], img_array[:, :, 1], img_array[:, :, 2]
            elif img.mode == 'RGBA':
                r, g, b, a = img_array[:, :, 0], img_array[:, :, 1], img_array[:, :, 2], img_array[:, :, 3]
            else:
                img_rgb = img.convert('RGB')
                r, g, b = np.array(img_rgb)[:, :, 0], np.array(img_rgb)[:, :, 1], np.array(img_rgb)[:, :, 2]

            return {
                'dominant_colors': self._find_dominant_colors(img_array),
                'color_variance': {
                    'red': float(np.var(r)),
                    'green': float(np.var(g)),
                    'blue': float(np.var(b))
                },
                'brightness_distribution': self._analyze_brightness_distribution(img_array),
                'saturation_analysis': self._analyze_saturation(img_array)
            }
        except:
            return {'error': 'Color analysis failed'}

    def _analyze_image_compression(self, img: Image.Image, file_path: Path) -> Dict:
        """Analyze image compression artifacts"""
        try:
            # Check for JPEG compression artifacts
            compression_info = {
                'has_compression_artifacts': False,
                'compression_quality': 'unknown',
                'blocking_artifacts': False
            }

            if hasattr(img, 'format') and img.format == 'JPEG':
                # Analyze for blocking artifacts
                img_array = np.array(img.convert('RGB'))
                compression_info.update(self._detect_jpeg_artifacts(img_array))

            return compression_info
        except:
            return {'error': 'Compression analysis failed'}

    def _analyze_image_quality(self, img: Image.Image) -> Dict:
        """Analyze image quality metrics"""
        try:
            img_gray = img.convert('L')
            img_array = np.array(img_gray)

            # Sharpness analysis
            from scipy.ndimage import sobel
            sobel_x = sobel(img_array, axis=0)
            sobel_y = sobel(img_array, axis=1)
            sharpness = np.sqrt(sobel_x**2 + sobel_y**2).mean()

            # Noise analysis
            noise_level = np.std(img_array)

            # Contrast analysis
            contrast = (np.max(img_array) - np.min(img_array)) / (np.max(img_array) + np.min(img_array))

            return {
                'sharpness': float(sharpness),
                'noise_level': float(noise_level),
                'contrast': float(contrast),
                'quality_score': self._calculate_image_quality_score(sharpness, noise_level, contrast)
            }
        except:
            return {'error': 'Quality analysis failed'}

    def _analyze_image_patterns(self, img: Image.Image) -> Dict:
        """Analyze image patterns and textures"""
        try:
            img_gray = np.array(img.convert('L'))

            # Edge detection
            edges = self._detect_edges(img_gray)

            # Pattern analysis
            patterns = {
                'edge_density': np.mean(edges > 0),
                'texture_complexity': self._calculate_texture_complexity(img_gray),
                'symmetry_score': self._calculate_symmetry(img_gray),
                'pattern_uniformity': self._calculate_pattern_uniformity(img_gray)
            }

            return patterns
        except:
            return {'error': 'Pattern analysis failed'}

    def _detect_edges(self, img_array: np.ndarray) -> np.ndarray:
        """Simple edge detection"""
        from scipy.ndimage import sobel
        sobel_x = sobel(img_array, axis=0)
        sobel_y = sobel(img_array, axis=1)
        return np.sqrt(sobel_x**2 + sobel_y**2)

    def _calculate_texture_complexity(self, img_array: np.ndarray) -> float:
        """Calculate texture complexity"""
        from scipy.ndimage import generic_filter
        def entropy_kernel(values):
            hist, _ = np.histogram(values, bins=16, range=(0, 255))
            hist = hist[hist > 0]
            return -np.sum(hist * np.log2(hist / np.sum(hist)))

        return generic_filter(img_array.astype(float), entropy_kernel, size=5).mean()

    def _calculate_symmetry(self, img_array: np.ndarray) -> float:
        """Calculate image symmetry"""
        h, w = img_array.shape
        left = img_array[:, :w//2]
        right = np.fliplr(img_array[:, w//2:])
        symmetry = 1 - np.mean(np.abs(left - right)) / 255
        return float(symmetry)

    def _calculate_pattern_uniformity(self, img_array: np.ndarray) -> float:
        """Calculate pattern uniformity"""
        from scipy.stats import entropy
        hist, _ = np.histogram(img_array, bins=32, range=(0, 255))
        hist = hist[hist > 0]
        return entropy(hist) / np.log2(len(hist)) if len(hist) > 1 else 0

    def _find_dominant_colors(self, img_array: np.ndarray, k: int = 5) -> List[Dict]:
        """Find dominant colors in image"""
        try:
            from sklearn.cluster import KMeans

            # Reshape for clustering
            pixels = img_array.reshape(-1, 3) if img_array.shape[-1] >= 3 else img_array.reshape(-1, 1)

            if pixels.shape[1] == 1:
                # Grayscale
                hist, bins = np.histogram(pixels, bins=32, range=(0, 255))
                dominant = []
                for i in range(min(k, len(hist))):
                    idx = np.argmax(hist)
                    hist[idx] = 0
                    dominant.append({
                        'color': [bins[idx], bins[idx], bins[idx]],
                        'percentage': hist[idx] / np.sum(hist)
                    })
                return dominant

            # Color image
            kmeans = KMeans(n_clusters=k, n_init=10)
            kmeans.fit(pixels)
            colors = kmeans.cluster_centers_
            labels = kmeans.labels_

            dominant = []
            for i in range(k):
                mask = labels == i
                percentage = np.sum(mask) / len(mask)
                dominant.append({
                    'color': colors[i].tolist(),
                    'percentage': float(percentage)
                })

            return sorted(dominant, key=lambda x: x['percentage'], reverse=True)

        except:
            return []

    def _analyze_brightness_distribution(self, img_array: np.ndarray) -> Dict:
        """Analyze brightness distribution"""
        try:
            if len(img_array.shape) == 3:
                brightness = np.mean(img_array, axis=2)
            else:
                brightness = img_array

            return {
                'mean_brightness': float(np.mean(brightness)),
                'brightness_std': float(np.std(brightness)),
                'brightness_histogram': np.histogram(brightness, bins=16, range=(0, 255))[0].tolist()
            }
        except:
            return {'error': 'Brightness analysis failed'}

    def _analyze_saturation(self, img_array: np.ndarray) -> Dict:
        """Analyze color saturation"""
        try:
            if len(img_array.shape) == 3 and img_array.shape[2] >= 3:
                hsv = cv2.cvtColor(img_array, cv2.COLOR_RGB2HSV) if img_array.shape[2] == 3 else None
                if hsv is not None:
                    saturation = hsv[:, :, 1]
                    return {
                        'mean_saturation': float(np.mean(saturation)),
                        'saturation_std': float(np.std(saturation)),
                        'saturation_range': [float(np.min(saturation)), float(np.max(saturation))]
                    }
            return {'mean_saturation': 0.5, 'saturation_std': 0.0, 'saturation_range': [0.0, 1.0]}
        except:
            return {'error': 'Saturation analysis failed'}

    def _detect_jpeg_artifacts(self, img_array: np.ndarray) -> Dict:
        """Detect JPEG compression artifacts"""
        try:
            # Check for 8x8 blocking artifacts
            h, w, _ = img_array.shape
            block_artifacts = 0
            total_blocks = 0

            for i in range(0, h-8, 8):
                for j in range(0, w-8, 8):
                    block = img_array[i:i+8, j:j+8]
                    # Check if block has uniform DCT-like patterns
                    if np.std(block) < 10:  # Very uniform blocks
                        block_artifacts += 1
                    total_blocks += 1

            compression_ratio = block_artifacts / total_blocks if total_blocks > 0 else 0

            return {
                'has_blocking_artifacts': compression_ratio > 0.1,
                'compression_artifacts_score': compression_ratio,
                'estimated_quality': max(1, min(100, 100 - compression_ratio * 1000))
            }
        except:
            return {'has_blocking_artifacts': False, 'compression_artifacts_score': 0.0, 'estimated_quality': 85}

    def _calculate_image_quality_score(self, sharpness: float, noise_level: float, contrast: float) -> float:
        """Calculate overall image quality score"""
        try:
            # Normalize metrics
            sharpness_score = min(sharpness / 50, 1.0)  # Higher is better
            noise_score = max(0, 1.0 - noise_level / 50)  # Lower noise is better
            contrast_score = min(contrast * 2, 1.0)  # Higher contrast is better

            # Weighted average
            return (sharpness_score * 0.4 + noise_score * 0.3 + contrast_score * 0.3)

        except:
            return 0.5

    def _get_color_depth(self, img: Image.Image) -> int:
        """Get color depth of image"""
        try:
            if img.mode == '1':
                return 1
            elif img.mode == 'L':
                return 8
            elif img.mode in ('P', 'RGB', 'YCBCR'):
                return 24
            elif img.mode == 'RGBA':
                return 32
            else:
                return 24  # Default
        except:
            return 24

    def _extract_image_metadata_file(self, file_path: Path) -> Dict:
        """Extract image metadata from file"""
        try:
            with open(file_path, 'rb') as f:
                tags = exifread.process_file(f)

            metadata = {}
            for tag, value in tags.items():
                if tag not in ('JPEGThumbnail', 'TIFFThumbnail'):
                    metadata[tag] = str(value)

            return metadata
        except:
            return {}

    def _extract_image_metadata(self, img: Image.Image) -> Dict:
        """Extract embedded metadata from PIL image"""
        try:
            metadata = {}
            if hasattr(img, '_getexif') and img._getexif():
                exif_data = img._getexif()
                for tag_id, value in exif_data.items():
                    tag = ExifTags.TAGS.get(tag_id, tag_id)
                    metadata[tag] = str(value)
            return metadata
        except:
            return {}

    # Video analysis helpers
    def _analyze_video_frames(self, clip: VideoFileClip) -> Dict:
        """Analyze video frames"""
        try:
            # Sample frames throughout video
            duration = clip.duration
            sample_times = np.linspace(0, duration, min(10, int(duration)))

            frame_analyses = []
            for t in sample_times:
                frame = clip.get_frame(t)
                # Basic frame analysis
                frame_analysis = {
                    'timestamp': t,
                    'brightness': float(np.mean(frame)),
                    'contrast': float(np.std(frame)),
                    'has_faces': False  # Would integrate face detection
                }
                frame_analyses.append(frame_analysis)

            return {
                'sampled_frames': len(frame_analyses),
                'frame_analyses': frame_analyses,
                'avg_brightness': np.mean([f['brightness'] for f in frame_analyses]),
                'avg_contrast': np.mean([f['contrast'] for f in frame_analyses])
            }
        except:
            return {'error': 'Frame analysis failed'}

    def _analyze_video_quality(self, clip: VideoFileClip) -> Dict:
        """Analyze video quality"""
        try:
            # Basic quality metrics
            return {
                'resolution': clip.size,
                'fps': clip.fps,
                'duration': clip.duration,
                'estimated_bitrate': 'unknown',  # Would need more analysis
                'compression_artifacts': False  # Would need detailed analysis
            }
        except:
            return {'error': 'Quality analysis failed'}

    def _analyze_video_content_features(self, clip: VideoFileClip) -> Dict:
        """Analyze video content features"""
        try:
            return {
                'motion_intensity': 'unknown',  # Would need optical flow analysis
                'scene_changes': 0,  # Would need scene detection
                'audio_sync': clip.audio is not None,
                'content_complexity': 'medium'  # Would need detailed analysis
            }
        except:
            return {'error': 'Content analysis failed'}

    def _analyze_audio_from_video(self, audio) -> Dict:
        """Analyze audio from video"""
        try:
            # This would extract audio features
            return {
                'duration': audio.duration if hasattr(audio, 'duration') else 0,
                'channels': 'unknown',
                'quality': 'unknown'
            }
        except:
            return {'error': 'Audio analysis failed'}

    # Audio analysis helpers
    def _analyze_audio_quality(self, y: np.ndarray, sr: int) -> Dict:
        """Analyze audio quality"""
        try:
            # Signal-to-noise ratio approximation
            signal_power = np.mean(y**2)
            noise_power = np.var(y) * 0.01  # Rough estimate
            snr = 10 * np.log10(signal_power / noise_power) if noise_power > 0 else 60

            # Dynamic range
            dynamic_range = np.max(y) - np.min(y)

            return {
                'snr_db': float(snr),
                'dynamic_range': float(dynamic_range),
                'sample_rate_quality': 'high' if sr >= 44100 else 'medium' if sr >= 22050 else 'low',
                'bit_depth_estimate': 'unknown'  # Would need more analysis
            }
        except:
            return {'error': 'Quality analysis failed'}

    def _analyze_audio_content(self, y: np.ndarray, sr: int) -> Dict:
        """Analyze audio content"""
        try:
            # Basic audio features using librosa
            mfccs = librosa.feature.mfcc(y=y, sr=sr, n_mfcc=13)
            chroma = librosa.feature.chroma_stft(y=y, sr=sr)
            spectral_centroid = librosa.feature.spectral_centroid(y=y, sr=sr)

            return {
                'mfcc_mean': np.mean(mfccs, axis=1).tolist(),
                'chroma_features': np.mean(chroma, axis=1).tolist(),
                'spectral_centroid': float(np.mean(spectral_centroid)),
                'zero_crossing_rate': float(np.mean(librosa.feature.zero_crossing_rate(y))),
                'tempo_estimate': float(librosa.beat.tempo(y=y, sr=sr)[0])
            }
        except:
            return {'error': 'Content analysis failed'}

    def _analyze_audio_technical(self, y: np.ndarray, sr: int) -> Dict:
        """Analyze technical audio properties"""
        try:
            # Frequency analysis
            fft = np.fft.fft(y)
            freqs = np.fft.fftfreq(len(y), 1/sr)
            magnitude = np.abs(fft)

            # Find dominant frequencies
            peak_freq_idx = np.argmax(magnitude[:len(magnitude)//2])
            dominant_freq = freqs[peak_freq_idx]

            return {
                'dominant_frequency': float(dominant_freq),
                'frequency_range': [float(np.min(freqs)), float(np.max(freqs))],
                'spectral_flatness': float(librosa.feature.spectral_flatness(y=y)[0].mean()),
                'rms_energy': float(librosa.feature.rms(y=y)[0].mean())
            }
        except:
            return {'error': 'Technical analysis failed'}

    # Document analysis helpers
    def _analyze_pptx_content(self, file_path: Path) -> Dict:
        """Analyze PowerPoint content"""
        try:
            prs = pptx.Presentation(str(file_path))

            analysis = {
                'slide_count': len(prs.slides),
                'properties': {
                    'author': prs.core_properties.author,
                    'created': prs.core_properties.created.isoformat() if prs.core_properties.created else None,
                    'title': prs.core_properties.title
                },
                'slides': []
            }

            for i, slide in enumerate(prs.slides):
                slide_analysis = {
                    'slide_number': i + 1,
                    'shape_count': len(slide.shapes),
                    'text_boxes': len([s for s in slide.shapes if hasattr(s, 'text')]),
                    'images': len([s for s in slide.shapes if hasattr(s, 'image')]),
                    'tables': len([s for s in slide.shapes if hasattr(s, 'table')])
                }
                analysis['slides'].append(slide_analysis)

            return analysis

        except Exception as e:
            return {'error': f'PPTX analysis failed: {e}'}

    def _analyze_xlsx_content(self, file_path: Path) -> Dict:
        """Analyze Excel content"""
        try:
            wb = openpyxl.load_workbook(str(file_path), read_only=True)

            analysis = {
                'worksheet_count': len(wb.worksheets),
                'properties': {
                    'creator': wb.properties.creator,
                    'created': wb.properties.created.isoformat() if wb.properties.created else None,
                    'title': wb.properties.title
                },
                'worksheets': []
            }

            for ws in wb.worksheets:
                ws_analysis = {
                    'name': ws.title,
                    'row_count': ws.max_row,
                    'column_count': ws.max_column,
                    'cell_count': 0,
                    'has_formulas': False
                }

                # Count non-empty cells (limited for performance)
                cell_count = 0
                has_formulas = False
                for row in ws.iter_rows(min_row=1, max_row=min(1000, ws.max_row)):
                    for cell in row:
                        if cell.value is not None:
                            cell_count += 1
                            if str(cell.value).startswith('='):
                                has_formulas = True

                ws_analysis['cell_count'] = cell_count
                ws_analysis['has_formulas'] = has_formulas
                analysis['worksheets'].append(ws_analysis)

            wb.close()
            return analysis

        except Exception as e:
            return {'error': f'XLSX analysis failed: {e}'}

    def _detect_text_language(self, text: str) -> str:
        """Detect text language (simplified)"""
        try:
            # Simple language detection based on common words
            text_lower = text.lower()

            if any(word in text_lower for word in ['the', 'and', 'is', 'in', 'to']):
                return 'english'
            elif any(word in text_lower for word in ['el', 'la', 'de', 'que', 'y']):
                return 'spanish'
            elif any(word in text_lower for word in ['der', 'die', 'das', 'und', 'ist']):
                return 'german'
            else:
                return 'unknown'
        except:
            return 'unknown'

    def _calculate_readability(self, text: str) -> Dict:
        """Calculate text readability metrics"""
        try:
            words = text.split()
            sentences = text.split('.')
            syllables = sum(self._count_syllables(word) for word in words)

            if len(sentences) == 0 or len(words) == 0:
                return {'flesch_score': 0, 'grade_level': 'unknown'}

            # Flesch Reading Ease
            flesch = 206.835 - 1.015 * (len(words) / len(sentences)) - 84.6 * (syllables / len(words))

            # Grade level
            grade_level = 0.39 * (len(words) / len(sentences)) + 11.8 * (syllables / len(words)) - 15.59

            return {
                'flesch_score': flesch,
                'grade_level': grade_level,
                'avg_words_per_sentence': len(words) / len(sentences),
                'avg_syllables_per_word': syllables / len(words)
            }
        except:
            return {'flesch_score': 0, 'grade_level': 'unknown'}

    def _count_syllables(self, word: str) -> int:
        """Count syllables in a word (simplified)"""
        word = word.lower()
        count = 0
        vowels = "aeiouy"
        if word[0] in vowels:
            count += 1
        for i in range(1, len(word)):
            if word[i] in vowels and word[i - 1] not in vowels:
                count += 1
        if word.endswith("e"):
            count -= 1
        if count == 0:
            count += 1
        return count

    def _analyze_text_sentiment(self, text: str) -> Dict:
        """Analyze text sentiment (simplified)"""
        try:
            positive_words = ['good', 'great', 'excellent', 'amazing', 'wonderful', 'fantastic', 'love', 'like', 'best']
            negative_words = ['bad', 'terrible', 'awful', 'hate', 'worst', 'horrible', 'disgusting', 'poor']

            words = text.lower().split()
            positive_count = sum(1 for word in words if word in positive_words)
            negative_count = sum(1 for word in words if word in negative_words)

            total_sentiment_words = positive_count + negative_count
            if total_sentiment_words == 0:
                sentiment = 'neutral'
                score = 0
            else:
                score = (positive_count - negative_count) / total_sentiment_words
                if score > 0.1:
                    sentiment = 'positive'
                elif score < -0.1:
                    sentiment = 'negative'
                else:
                    sentiment = 'neutral'

            return {
                'sentiment': sentiment,
                'score': score,
                'positive_words': positive_count,
                'negative_words': negative_count
            }
        except:
            return {'sentiment': 'neutral', 'score': 0}

    # Metadata extraction helpers
    def _extract_video_metadata(self, file_path: Path) -> Dict:
        """Extract video metadata"""
        try:
            clip = VideoFileClip(str(file_path))
            metadata = {
                'duration': clip.duration,
                'fps': clip.fps,
                'size': clip.size,
                'audio_codec': 'unknown',
                'video_codec': 'unknown'
            }
            clip.close()
            return metadata
        except:
            return {}

    def _extract_audio_metadata(self, file_path: Path) -> Dict:
        """Extract audio metadata"""
        try:
            y, sr = librosa.load(str(file_path), sr=None)
            return {
                'duration': len(y) / sr,
                'sample_rate': sr,
                'channels': 1 if y.ndim == 1 else y.shape[0]
            }
        except:
            return {}

    def _extract_document_metadata(self, file_path: Path) -> Dict:
        """Extract document metadata"""
        try:
            ext = file_path.suffix.lower()
            if ext == '.pdf':
                doc = fitz.open(str(file_path))
                metadata = dict(doc.metadata)
                doc.close()
                return metadata
            elif ext == '.docx':
                doc = docx.Document(str(file_path))
                return {
                    'author': doc.core_properties.author,
                    'created': doc.core_properties.created.isoformat() if doc.core_properties.created else None,
                    'modified': doc.core_properties.modified.isoformat() if doc.core_properties.modified else None,
                    'title': doc.core_properties.title
                }
            else:
                return {}
        except:
            return {}

    def _get_file_system_metadata(self, file_path: Path) -> Dict:
        """Get file system metadata"""
        try:
            stat = file_path.stat()
            return {
                'size_bytes': stat.st_size,
                'created': datetime.fromtimestamp(stat.st_ctime).isoformat(),
                'modified': datetime.fromtimestamp(stat.st_mtime).isoformat(),
                'accessed': datetime.fromtimestamp(stat.st_atime).isoformat(),
                'permissions': oct(stat.st_mode)[-3:]
            }
        except:
            return {}

    def _calculate_file_hashes(self, file_path: Path) -> Dict:
        """Calculate file hashes for integrity checking"""
        try:
            hash_md5 = hashlib.md5()
            hash_sha256 = hashlib.sha256()

            with open(file_path, 'rb') as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
                    hash_sha256.update(chunk)

            return {
                'md5': hash_md5.hexdigest(),
                'sha256': hash_sha256.hexdigest()
            }
        except:
            return {}

    # AI vs Human analysis helpers
    def _analyze_image_ai_vs_human(self, content_analysis: Dict) -> Dict:
        """Analyze if image was likely created by AI or human"""
        indicators = []
        ai_score = 0.5

        try:
            # Check for common AI generation artifacts
            if content_analysis.get('compression_analysis', {}).get('has_blocking_artifacts'):
                indicators.append('JPEG blocking artifacts detected')
                ai_score += 0.2

            if content_analysis.get('quality_analysis', {}).get('quality_score', 0.5) < 0.3:
                indicators.append('Unusually low quality for modern image')
                ai_score += 0.1

            # Check color patterns
            color_analysis = content_analysis.get('color_analysis', {})
            if color_analysis.get('color_variance', {}).get('red', 0) > 10000:
                indicators.append('Unnatural color variance patterns')
                ai_score += 0.15

            # Check pattern uniformity
            pattern_analysis = content_analysis.get('pattern_analysis', {})
            if pattern_analysis.get('pattern_uniformity', 0.5) > 0.8:
                indicators.append('Highly uniform patterns (AI characteristic)')
                ai_score += 0.2

            # Check symmetry (AI often generates symmetric content)
            if pattern_analysis.get('symmetry_score', 0.5) > 0.8:
                indicators.append('Unusually high symmetry')
                ai_score += 0.1

            # Human indicators
            if content_analysis.get('metadata', {}).get('embedded_metadata'):
                indicators.append('Contains camera metadata (human-created)')
                ai_score -= 0.2

            if content_analysis.get('quality_analysis', {}).get('noise_level', 0) > 20:
                indicators.append('Natural noise patterns (human photography)')
                ai_score -= 0.1

        except Exception as e:
            indicators.append(f'Analysis error: {e}')

        return {
            'ai_probability': min(1.0, max(0.0, ai_score)),
            'human_probability': 1.0 - min(1.0, max(0.0, ai_score)),
            'indicators': indicators
        }

    def _analyze_video_ai_vs_human(self, content_analysis: Dict) -> Dict:
        """Analyze if video was likely created by AI or human"""
        indicators = []
        ai_score = 0.5

        try:
            # Check frame consistency
            frame_analysis = content_analysis.get('frame_analysis', {})
            if frame_analysis.get('avg_brightness', 0.5) > 0.9 or frame_analysis.get('avg_brightness', 0.5) < 0.1:
                indicators.append('Unnatural brightness levels')
                ai_score += 0.1

            # Check for compression artifacts
            quality_analysis = content_analysis.get('quality_analysis', {})
            if quality_analysis.get('compression_artifacts'):
                indicators.append('Heavy compression artifacts')
                ai_score += 0.15

            # Human indicators
            if content_analysis.get('audio_present'):
                indicators.append('Contains audio track (likely human-recorded)')
                ai_score -= 0.2

            if content_analysis.get('frame_count', 0) > 1000:
                indicators.append('Long duration video (likely human-recorded)')
                ai_score -= 0.1

        except Exception as e:
            indicators.append(f'Analysis error: {e}')

        return {
            'ai_probability': min(1.0, max(0.0, ai_score)),
            'human_probability': 1.0 - min(1.0, max(0.0, ai_score)),
            'indicators': indicators
        }

    def _analyze_audio_ai_vs_human(self, content_analysis: Dict) -> Dict:
        """Analyze if audio was likely created by AI or human"""
        indicators = []
        ai_score = 0.5

        try:
            # Check quality metrics
            quality_analysis = content_analysis.get('quality_analysis', {})
            if quality_analysis.get('snr_db', 30) > 50:
                indicators.append('Unusually high signal-to-noise ratio')
                ai_score += 0.2

            # Check technical features
            technical_analysis = content_analysis.get('technical_analysis', {})
            if technical_analysis.get('spectral_flatness', 0.5) > 0.8:
                indicators.append('Unnatural spectral characteristics')
                ai_score += 0.15

            # Human indicators
            if quality_analysis.get('dynamic_range', 0) > 0.8:
                indicators.append('Wide dynamic range (human recording)')
                ai_score -= 0.2

            content_features = content_analysis.get('content_analysis', {})
            if content_features.get('zero_crossing_rate', 0) > 0.1:
                indicators.append('Natural audio characteristics')
                ai_score -= 0.1

        except Exception as e:
            indicators.append(f'Analysis error: {e}')

        return {
            'ai_probability': min(1.0, max(0.0, ai_score)),
            'human_probability': 1.0 - min(1.0, max(0.0, ai_score)),
            'indicators': indicators
        }

    def _analyze_document_ai_vs_human(self, content_analysis: Dict) -> Dict:
        """Analyze if document was likely created by AI or human"""
        indicators = []
        ai_score = 0.5

        try:
            if content_analysis.get('content_type') == 'document':
                # Check text analysis
                if 'text_analysis' in content_analysis:
                    text_analysis = content_analysis['text_analysis']

                    # Check readability
                    readability = text_analysis.get('readability_score', {})
                    if readability.get('flesch_score', 60) > 90:
                        indicators.append('Unusually high readability score')
                        ai_score += 0.1

                    # Check language patterns
                    if text_analysis.get('language_detection') == 'unknown':
                        indicators.append('Unusual language patterns')
                        ai_score += 0.15

                    # Human indicators
                    if readability.get('grade_level', 8) > 12:
                        indicators.append('Complex language patterns (human writing)')
                        ai_score -= 0.1

            # Check metadata
            if content_analysis.get('metadata'):
                indicators.append('Contains creation metadata (likely human-created)')
                ai_score -= 0.2

        except Exception as e:
            indicators.append(f'Analysis error: {e}')

        return {
            'ai_probability': min(1.0, max(0.0, ai_score)),
            'human_probability': 1.0 - min(1.0, max(0.0, ai_score)),
            'indicators': indicators
        }

    # Summary generation helpers
    def _summarize_content(self, content_analysis: Dict) -> str:
        """Generate content summary"""
        try:
            if 'error' in content_analysis:
                return f"Analysis failed: {content_analysis['error']}"

            if content_analysis.get('content_type') == 'image':
                dims = content_analysis.get('dimensions', 'Unknown')
                mode = content_analysis.get('mode', 'Unknown')
                return f"Image: {dims[0]}x{dims[1]} pixels, {mode} mode"

            elif content_analysis.get('content_type') == 'video':
                duration = content_analysis.get('duration', 0)
                fps = content_analysis.get('fps', 0)
                size = content_analysis.get('size', (0, 0))
                return f"Video: {size[0]}x{size[1]}, {duration:.1f}s, {fps}fps"

            elif content_analysis.get('content_type') == 'audio':
                duration = content_analysis.get('duration', 0)
                sr = content_analysis.get('sample_rate', 0)
                channels = content_analysis.get('channels', 1)
                return f"Audio: {duration:.1f}s, {sr}Hz, {channels} channel{'s' if channels != 1 else ''}"

            elif content_analysis.get('content_type') == 'document':
                if 'page_count' in content_analysis:
                    pages = content_analysis['page_count']
                    return f"PDF Document: {pages} pages"
                elif 'paragraph_count' in content_analysis:
                    paras = content_analysis['paragraph_count']
                    return f"Word Document: {paras} paragraphs"
                else:
                    return "Document file"

            else:
                return "Unsupported content type"

        except:
            return "Content summary unavailable"

    def _extract_key_findings(self, content_analysis: Dict, metadata: Dict, ai_human_analysis: Dict) -> List[str]:
        """Extract key findings from analysis"""
        findings = []

        try:
            # Content findings
            if content_analysis.get('content_type') == 'image':
                quality = content_analysis.get('quality_analysis', {})
                if quality.get('quality_score', 0.5) < 0.3:
                    findings.append("Low quality image detected")
                if quality.get('noise_level', 0) > 30:
                    findings.append("High noise levels detected")

            # AI/Human findings
            ai_prob = ai_human_analysis.get('ai_probability', 0.5)
            if ai_prob > 0.7:
                findings.append("High likelihood of AI-generated content")
            elif ai_prob < 0.3:
                findings.append("Likely human-created content")

            # Metadata findings
            if metadata.get('file_system', {}).get('size_bytes', 0) > 100 * 1024 * 1024:  # 100MB
                findings.append("Large file size detected")

        except:
            findings.append("Key findings extraction failed")

        return findings

    def _assess_risks(self, ai_human_analysis: Dict) -> Dict:
        """Assess risks based on analysis"""
        try:
            ai_prob = ai_human_analysis.get('ai_probability', 0.5)

            if ai_prob > 0.8:
                risk_level = 'high'
                description = 'Strong indicators of AI-generated content'
            elif ai_prob > 0.6:
                risk_level = 'medium'
                description = 'Moderate likelihood of AI-generated content'
            elif ai_prob < 0.4:
                risk_level = 'low'
                description = 'Likely authentic human-created content'
            else:
                risk_level = 'unknown'
                description = 'Unable to determine content origin'

            return {
                'risk_level': risk_level,
                'description': description,
                'confidence': ai_human_analysis.get('confidence', 'low')
            }

        except:
            return {'risk_level': 'unknown', 'description': 'Risk assessment failed'}

    def _generate_recommendations(self, ai_human_analysis: Dict) -> List[str]:
        """Generate recommendations based on analysis"""
        recommendations = []

        try:
            ai_prob = ai_human_analysis.get('ai_probability', 0.5)

            if ai_prob > 0.7:
                recommendations.extend([
                    'Verify content authenticity through multiple sources',
                    'Check for digital watermarking or signatures',
                    'Consider additional forensic analysis',
                    'Be cautious with content attribution'
                ])
            elif ai_prob < 0.3:
                recommendations.extend([
                    'Content appears authentic',
                    'Standard verification procedures recommended',
                    'Check metadata for additional context'
                ])
            else:
                recommendations.extend([
                    'Further analysis recommended',
                    'Consider expert verification',
                    'Review content in context of other evidence'
                ])

        except:
            recommendations.append('Unable to generate recommendations')

        return recommendations

    def _format_file_size(self, size_bytes: int) -> str:
        """Format file size in human readable format"""
        for unit in ['B', 'KB', 'MB', 'GB', 'TB']:
            if size_bytes < 1024.0:
                return f"{size_bytes:.1f} {unit}"
            size_bytes /= 1024.0
        return f"{size_bytes:.1f} PB"

    def _analyze_generic_content(self, file_path: Path) -> Dict:
        """Analyze generic file content"""
        try:
            # Read first 1KB for basic analysis
            with open(file_path, 'rb') as f:
                header = f.read(1024)

            return {
                'file_header': header[:50].hex(),
                'is_binary': self._is_binary_file(header),
                'entropy': self._calculate_entropy(header),
                'content_type': 'generic'
            }
        except:
            return {'error': 'Generic analysis failed'}

    def _analyze_archive_content(self, file_path: Path) -> Dict:
        """Analyze archive content"""
        try:
            # Basic archive analysis
            return {
                'archive_type': file_path.suffix.upper(),
                'estimated_files': 'unknown',  # Would need to extract
                'compression_method': 'unknown',
                'content_type': 'archive'
            }
        except:
            return {'error': 'Archive analysis failed'}

    def _is_binary_file(self, data: bytes) -> bool:
        """Check if file is binary"""
        text_chars = bytearray({7,8,9,10,12,13,27} | set(range(0x20, 0x100)) - {0x7f})
        return bool(data.translate(None, text_chars))

    def _calculate_entropy(self, data: bytes) -> float:
        """Calculate Shannon entropy of data"""
        if not data:
            return 0.0

        entropy = 0.0
        for i in range(256):
            p = data.count(i) / len(data)
            if p > 0:
                entropy -= p * np.log2(p)

        return entropy